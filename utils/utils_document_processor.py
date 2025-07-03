import os
from typing import List, Dict, Tuple
import PyPDF2
from docx import Document
import pandas as pd
from pptx import Presentation
import logging
import sys
import os

# Add the project root to the Python path for standalone execution
current_dir_dp = os.path.dirname(os.path.abspath(__file__))
project_root_dp = os.path.abspath(os.path.join(current_dir_dp, '..', '..'))
if project_root_dp not in sys.path:
    sys.path.insert(0, project_root_dp)

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

class DocumentProcessor:
    """
    Handles parsing and extracting text from various document types.
    """
    def __init__(self):
        pass

    def parse_document(self, file_path: str) -> Tuple[str, str]:
        """
        Parses a document based on its file extension and extracts all text.

        Args:
            file_path (str): The path to the document file.

        Returns:
            Tuple[str, str]: A tuple containing (extracted_text, document_type).
                             Returns an empty string and 'unknown' type if parsing fails.
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        document_type = file_extension[1:] # Remove the dot

        try:
            if file_extension == '.pdf':
                return self._parse_pdf(file_path), document_type
            elif file_extension == '.docx':
                return self._parse_docx(file_path), document_type
            elif file_extension == '.csv':
                return self._parse_csv(file_path), document_type
            elif file_extension == '.pptx':
                return self._parse_pptx(file_path), document_type
            elif file_extension in ['.txt', '.md']:
                return self._parse_text_file(file_path), document_type
            else:
                logging.warning(f"Unsupported file type: {file_extension}")
                return "", "unsupported"
        except Exception as e:
            logging.error(f"Error parsing {file_path}: {e}", exc_info=True)
            return "", "error"

    def _parse_pdf(self, file_path: str) -> str:
        """Extracts text from a PDF file."""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    text += page.extract_text() or ""
            return text
        except Exception as e:
            logging.error(f"Failed to parse PDF {file_path}: {e}")
            return ""

    def _parse_docx(self, file_path: str) -> str:
        """Extracts text from a DOCX file."""
        text = ""
        try:
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
            return text
        except Exception as e:
            logging.error(f"Failed to parse DOCX {file_path}: {e}")
            return ""

    def _parse_csv(self, file_path: str) -> str:
        """Extracts text from a CSV file."""
        try:
            df = pd.read_csv(file_path)
            # Convert DataFrame to a string representation for ingestion
            return df.to_string(index=False)
        except Exception as e:
            logging.error(f"Failed to parse CSV {file_path}: {e}")
            return ""

    def _parse_pptx(self, file_path: str) -> str:
        """Extracts text from a PPTX file."""
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text += shape.text_frame.text + "\n"
            return text
        except Exception as e:
            logging.error(f"Failed to parse PPTX {file_path}: {e}")
            return ""

    def _parse_text_file(self, file_path: str) -> str:
        """Extracts text from a TXT or Markdown file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logging.error(f"Failed to parse text file {file_path}: {e}")
            return ""

    def split_text_into_chunks(self,
                               text: str,
                               chunk_size: int = 1000,
                               chunk_overlap: int = 100) -> List[str]:
        """
        Splits a given text into overlapping chunks.

        Args:
            text (str): The input text to split.
            chunk_size (int): The maximum size of each chunk.
            chunk_overlap (int): The number of characters to overlap between chunks.

        Returns:
            List[str]: A list of text chunks.
        """
        if not text:
            return []

        chunks = []
        start_index = 0
        while start_index < len(text):
            end_index = min(start_index + chunk_size, len(text))
            chunk = text[start_index:end_index]
            chunks.append(chunk)
            if end_index == len(text):
                break
            start_index += (chunk_size - chunk_overlap)
            # Ensure start_index doesn't go negative or past text length
            start_index = max(0, min(start_index, len(text) - 1)) # Adjust for last chunk overlap

        return chunks

# Example Usage (for testing)
if __name__ == "__main__":
    import tempfile

    processor = DocumentProcessor()

    # Test TXT
    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp:
        tmp.write(b"This is a sample text file for testing. It has multiple sentences.")
        txt_path = tmp.name
    text, doc_type = processor.parse_document(txt_path)
    print(f"TXT Content: '{text[:50]}...', Type: {doc_type}")
    chunks = processor.split_text_into_chunks(text, chunk_size=30, chunk_overlap=5)
    print(f"TXT Chunks: {chunks}")
    os.remove(txt_path)

    # Test DOCX (requires a real docx file or creation)
    try:
        from docx import Document
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
            doc = Document()
            doc.add_paragraph("Hello from Word document.")
            doc.add_paragraph("This is the second paragraph.")
            doc.save(tmp.name)
            docx_path = tmp.name
        text, doc_type = processor.parse_document(docx_path)
        print(f"DOCX Content: '{text[:50]}...', Type: {doc_type}")
        os.remove(docx_path)
    except ImportError:
        print("python-docx not installed, skipping DOCX test.")
    except Exception as e:
        print(f"Error testing DOCX: {e}")

    # Test PDF (requires a dummy PDF or a library to create one)
    # print("\nSkipping PDF test as it requires a sample PDF file.")